import sys
sys.stdout.reconfigure(encoding='utf-8')
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Цветова схема (оранжево-тъмно като учебника) ──────────────────────────
C_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # почти черно
C_ORANGE = RGBColor(0xD6, 0x4A, 0x1A)   # оранжево (като учебника)
C_ACCENT = RGBColor(0xFF, 0x9A, 0x3C)   # светло оранжево
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BG     = RGBColor(0xFA, 0xF7, 0xF4)   # топъл бял фон
C_GRAY   = RGBColor(0x33, 0x33, 0x33)
C_BLUE   = RGBColor(0x1A, 0x5F, 0x9E)   # синьо за код
C_GREEN  = RGBColor(0x00, 0xAA, 0x44)   # зелено за верен отговор
C_CODE_BG = RGBColor(0xF0, 0xF0, 0xF8)  # светъл фон за код
C_KW     = RGBColor(0x00, 0x00, 0xCC)   # ключови думи

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def new_slide():
    return prs.slides.add_slide(BLANK)

def rect(sl, l, t, w, h, fill=None, line_color=None, line_pt=1.5):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color; s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    return s

def txb(sl, text, l, t, w, h, size=22, bold=False, color=C_GRAY,
        align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return tb

def mtxb(sl, lines, l, t, w, h, size=20, color=C_GRAY, spacing=6):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, bld, clr) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(spacing)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bld
        r.font.color.rgb = clr if clr else color
    return tb

def header(sl, title, sub=''):
    rect(sl, 0, 0, 13.33, 1.55, fill=C_DARK)
    rect(sl, 0, 1.55, 13.33, 0.08, fill=C_ORANGE)
    txb(sl, title, 0.35, 0.1, 12.6, 0.88, size=34, bold=True, color=C_WHITE)
    if sub:
        txb(sl, sub, 0.35, 1.0, 12.6, 0.52, size=19, color=C_ACCENT)

def code_box(sl, code, l, t, w, h, size=15):
    rect(sl, l, t, w, h, fill=C_CODE_BG, line_color=C_BLUE)
    txb(sl, code, l+0.18, t+0.15, w-0.36, h-0.3, size=size, color=C_DARK)

def add_notes(sl, text):
    notes_slide = sl.notes_slide
    notes_slide.notes_text_frame.text = text

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 1 — Заглавен
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=C_DARK)
rect(sl, 0, 0, 0.18, 7.5, fill=C_ORANGE)
rect(sl, 0, 7.35, 13.33, 0.15, fill=C_ORANGE)

txb(sl, 'Модул 2  ·  Урок 10  ·  Едномерен масив', 0.4, 0.6, 12, 0.6,
    size=20, color=C_ACCENT)
txb(sl, 'СОРТИРАНЕ НА МАСИВ', 0.4, 1.3, 12.5, 1.4,
    size=52, bold=True, color=C_WHITE)
rect(sl, 0.4, 2.85, 5, 0.07, fill=C_ORANGE)

mtxb(sl, [
    ('▸  Алгоритми за сортиране', False, C_WHITE),
    ('▸  Bubble Sort и Selection Sort', False, C_WHITE),
    ('▸  Сортиране на низове и обекти', False, C_WHITE),
    ('▸  Array.Sort() — вградени методи', False, C_WHITE),
], 0.4, 3.05, 9, 2.8, size=24, color=C_WHITE, spacing=12)

txb(sl, '11Б клас  ·  C# програмиране', 0.4, 6.75, 10, 0.5,
    size=17, color=RGBColor(0x80,0x80,0x90))

add_notes(sl, """ЗАГЛАВЕН СЛАЙД — Урок 10: Сортиране на масив

Цел на урока:
Учениците да разберат:
1. Защо е нужно сортирането
2. Как работят два основни алгоритъма: Bubble Sort и Selection Sort
3. Как се сортират числа, низове и обекти в C#
4. Как да използват вградения Array.Sort()

Въвеждащ въпрос:
"Представете си, че имате списък от 1000 ученика и трябва да намерите Иван Иванов. Ще е по-лесно ако списъкът е сортиран по азбучен ред, нали?"
→ Ето затова сортирането е толкова важно!

Предварителни знания, от които се нуждаят:
- Деклариране и обхождане на масив
- for и foreach цикли
- if условия""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 2 — Какво е сортиране?
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=C_BG)
header(sl, '📋  Какво е сортиране на масив?', 'Подреждане на елементите по даден критерий')

rect(sl, 0.35, 1.75, 12.6, 1.55, fill=C_DARK)
txb(sl, '  Сортирането на масив = подреждане на елементите му по даден критерий.',
    0.45, 1.82, 12.3, 0.65, size=24, bold=True, color=C_WHITE)
txb(sl, '  Най-честият критерий е стойността на елементите.',
    0.45, 2.48, 12.3, 0.65, size=22, color=C_ACCENT)

rect(sl, 0.35, 3.45, 5.9, 3.75, fill=C_WHITE, line_color=C_ORANGE)
txb(sl, '🔢  Числов масив', 0.55, 3.6, 5.5, 0.5, size=23, bold=True, color=C_ORANGE)
mtxb(sl, [
    ('▸  Възходящ (нарастващ) ред:', False, C_GRAY),
    ('    1, 2, 3, 4, 5, 9', False, C_BLUE),
    ('', False, C_GRAY),
    ('▸  Низходящ (намаляващ) ред:', False, C_GRAY),
    ('    9, 5, 4, 3, 2, 1', False, C_BLUE),
], 0.55, 4.15, 5.5, 2.85, size=22, color=C_GRAY, spacing=10)

rect(sl, 6.6, 3.45, 6.4, 3.75, fill=C_WHITE, line_color=C_BLUE)
txb(sl, '🔤  Масив от низове', 6.8, 3.6, 6.0, 0.5, size=23, bold=True, color=C_BLUE)
mtxb(sl, [
    ('▸  Лексикографски (азбучен) ред:', False, C_GRAY),
    ('    "Ани", "Боби", "Мими"', False, C_BLUE),
    ('', False, C_GRAY),
    ('▸  По дължина на думата', False, C_GRAY),
    ('▸  По друг потребителски критерий', False, C_GRAY),
], 6.8, 4.15, 6.0, 2.85, size=22, color=C_GRAY, spacing=10)

add_notes(sl, """КАКВО Е СОРТИРАНЕ?

Обяснение за класа:
Представете си, че имате карти с числа: 5, 2, 9, 7, 4, 1
"Как бихте ги наредили от най-малкото към най-голямото?"
→ Учениците интуитивно знаят как, но компютърът трябва да бъде НАУЧЕН!

Важни термини:
- Възходящ ред = нарастващ = ascending = от малко към голямо
- Низходящ ред = намаляващ = descending = от голямо към малко
- Лексикографски ред = азбучен ред (буква по буква)

Пример за лексикографски ред:
"apple" < "banana" (защото 'a' < 'b')
"app" < "apple" (по-късата дума е по-малка при еднакво начало)

Реален пример: телефонен указател, списък с ученици, класация в игра.""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 3 — Видове алгоритми
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=C_BG)
header(sl, '⚙️  Алгоритми за сортиране', 'Бавни и бързи методи')

txb(sl, 'Всички алгоритми дават ЕДИН И СЪЩ резултат — разликата е в скоростта!',
    0.35, 1.75, 12.6, 0.55, size=22, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

rect(sl, 0.35, 2.45, 5.95, 4.7, fill=C_WHITE, line_color=C_ORANGE)
rect(sl, 0.35, 2.45, 5.95, 0.58, fill=C_ORANGE)
txb(sl, '  🐌  БАВНИ алгоритми', 0.45, 2.53, 5.7, 0.42, size=22, bold=True, color=C_WHITE)
mtxb(sl, [
    ('Подходящи за малък брой елементи', True, C_ORANGE),
    ('', False, C_GRAY),
    ('▸  Метод на мехурчето (Bubble Sort)', False, None),
    ('▸  Метод на пряката селекция (Selection Sort)', False, None),
    ('▸  Сортиране чрез вмъкване (Insertion Sort)', False, None),
    ('', False, C_GRAY),
    ('Сложност:  O(n²)', True, C_ORANGE),
    ('→ 100 елемента = 10 000 сравнения', False, None),
], 0.55, 3.1, 5.6, 3.85, size=21, color=C_GRAY, spacing=9)

rect(sl, 6.65, 2.45, 6.3, 4.7, fill=C_WHITE, line_color=C_BLUE)
rect(sl, 6.65, 2.45, 6.3, 0.58, fill=C_BLUE)
txb(sl, '  🚀  БЪРЗИ алгоритми', 6.75, 2.53, 6.1, 0.42, size=22, bold=True, color=C_WHITE)
mtxb(sl, [
    ('Ефективни при голям брой елементи', True, C_BLUE),
    ('', False, C_GRAY),
    ('▸  Бързо сортиране (Quick Sort)', False, None),
    ('▸  Сортиране чрез сливане (Merge Sort)', False, None),
    ('▸  Пирамидално (Heap Sort)', False, None),
    ('', False, C_GRAY),
    ('Сложност:  O(n log n)', True, C_BLUE),
    ('→ 100 елемента ≈ 664 сравнения', False, None),
], 6.85, 3.1, 6.0, 3.85, size=21, color=C_GRAY, spacing=9)

add_notes(sl, """ВИДОВЕ АЛГОРИТМИ ЗА СОРТИРАНЕ

Обяснение на O(n²) и O(n log n) — просто!:
"Ако имате 10 елемента:
- Бавен алгоритъм: ~100 операции (10²)
- Бърз алгоритъм: ~33 операции (10 × log₂10 ≈ 33)

Ако имате 1 000 000 елемента:
- Бавен: 1 000 000 000 000 операции (1 трилион!) → часове
- Бърз: ~20 000 000 операции → секунди"

В този урок учим бавните — те са по-лесни за разбиране!
Бързите алгоритми са в по-напреднал курс.

Визуализация: Покажете visualgo.net/en/sorting на проектора!
- Изберете Bubble Sort и натиснете Sort
- Учениците виждат анимация на алгоритъма
- После сравнете с Quick Sort — видима разлика в скоростта!""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 4 — Bubble Sort: идеята
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=C_BG)
header(sl, '🫧  Метод на мехурчето (Bubble Sort)', 'Сравняваме съседни елементи и ги разменяме')

txb(sl, 'Идея: Сравняваме два съседни елемента. Ако не са наредени — разменяме ги.',
    0.35, 1.75, 12.6, 0.55, size=23, bold=True, color=C_DARK)

# Визуализация: масив 5 2 9 7 4 1 → един обход
steps = [
    ([5,2,9,7,4,1], 0, 1, True,  'Сравняваме 5 и 2 → 5>2 → разменяме'),
    ([2,5,9,7,4,1], 1, 2, False, 'Сравняваме 5 и 9 → 5<9 → оставяме'),
    ([2,5,9,7,4,1], 2, 3, True,  'Сравняваме 9 и 7 → 9>7 → разменяме'),
    ([2,5,7,9,4,1], 3, 4, True,  'Сравняваме 9 и 4 → 9>4 → разменяме'),
    ([2,5,7,4,9,1], 4, 5, True,  'Сравняваме 9 и 1 → 9>1 → разменяме → 9 е на място! 🎯'),
]

cell_w = 1.15
x_start = 0.5
y_start = 2.5

for row_i, (arr, j, j1, swapped, desc) in enumerate(steps):
    y = y_start + row_i * 0.88
    color_row = RGBColor(0xFF,0xEB,0xCC) if swapped else RGBColor(0xE8,0xF5,0xE9)
    for ci, val in enumerate(arr):
        is_active = ci == j or ci == j1
        bg = C_ORANGE if is_active else RGBColor(0xEE,0xEE,0xEE)
        fg = C_WHITE if is_active else C_DARK
        rect(sl, x_start + ci * cell_w, y, cell_w - 0.05, 0.72, fill=bg)
        txb(sl, str(val), x_start + ci * cell_w + 0.35, y + 0.12,
            0.45, 0.5, size=22, bold=True, color=fg, align=PP_ALIGN.CENTER)
    icon = '🔄' if swapped else '✓'
    txb(sl, f'{icon}  {desc}', x_start + 6.2 * cell_w + 0.1, y + 0.12,
        5.8, 0.6, size=17, color=C_DARK)

txb(sl, 'След 1 обход: най-голямото число "изплува" на последна позиция!',
    0.35, 7.0, 12.6, 0.42, size=20, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

add_notes(sl, """BUBBLE SORT — ИДЕЯ И ВИЗУАЛИЗАЦИЯ

Обяснение стъпка по стъпка:
Покажете физически с карти на дъската!
- Вземете 6 карти с числата 5, 2, 9, 7, 4, 1
- Обходете ги двойка по двойка и разменяйте ако трябва
- След 1 обход, 9 е на последно място

Защо се казва "мехурче"?
"Представете си, че всяко голямо число е мехурче — то се "издига" нагоре (към края на масива) с всеки обход, точно като мехурчета в газирана вода!"

Колко обхода са нужни?
- n-1 обхода за n елемента
- При всеки следващ обход проверяваме с 1 елемент по-малко (последните са вече наредени)
- Оптимизация: ако в даден обход няма размяна → масивът е вече сортиран!

Важно за кода: Разменяме с ВРЕМЕННА ПРОМЕНЛИВА temp:
int temp = a[j];     ← запомняме a[j]
a[j] = a[j+1];      ← презаписваме a[j]
a[j+1] = temp;      ← слагаме старото a[j] на новото място""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 5 — Bubble Sort: код
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8,0xF8,0xFF))
header(sl, '{ }  Bubble Sort — Код в C#', 'Реализация на метода на мехурчето')

bubble_code = """static void BubbleSort(int[] numbers)
{
    for (int i = 0; i < numbers.Length - 1; i++)
    {
        for (int j = 0; j < numbers.Length - 1 - i; j++)
        {
            if (numbers[j] > numbers[j + 1])
            {
                int temp      = numbers[j];
                numbers[j]    = numbers[j + 1];
                numbers[j + 1] = temp;
            }
        }
    }
}"""

main_code = """static void Main(string[] args)
{
    int[] numbers = Console.ReadLine()
        .Split()
        .Select(int.Parse)
        .ToArray();

    BubbleSort(numbers);

    Console.WriteLine(
        string.Join(" ", numbers));
}"""

code_box(sl, bubble_code, 0.35, 1.75, 7.5, 5.3, size=16)
code_box(sl, main_code,   8.05, 1.75, 5.0, 5.3, size=16)

add_notes(sl, """BUBBLE SORT — КОД

Обяснение на кода ред по ред:

Външният цикъл: for (int i = 0; i < numbers.Length - 1; i++)
→ Брои обходите. Правим n-1 обхода.

Вътрешният цикъл: for (int j = 0; j < numbers.Length - 1 - i; j++)
→ Обхожда двойките. Спира по-рано (- i) защото последните i елемента вече са наредени!

Условието: if (numbers[j] > numbers[j + 1])
→ Ако текущият е ПО-ГОЛЯМ от следващия → трябва размяна
→ За НАМАЛЯВАЩ ред: сменете > с <

Размяната с temp — ЗАДЪЛЖИТЕЛНО!
int temp = numbers[j];       // запомняме
numbers[j] = numbers[j+1];  // презаписваме
numbers[j+1] = temp;         // слагаме стария

БЕЗ temp: numbers[j] = numbers[j+1]; → ГУБИМ стойността на numbers[j]!

Main метода:
.Split() → разделя входния ред по интервали
.Select(int.Parse) → преобразува всяка дума в число
.ToArray() → събира в масив

Пример: вход "5 2 9 7 4 1" → изход "1 2 4 5 7 9" """)

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 6 — Selection Sort: идеята
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=C_BG)
header(sl, '🎯  Метод на пряката селекция (Selection Sort)', 'Намираме минималния елемент и го поставяме на място')

txb(sl, 'Идея: При всеки обход намираме най-малкия елемент и го разменяме с текущата позиция.',
    0.35, 1.75, 12.6, 0.55, size=22, bold=True, color=C_DARK)

steps_sel = [
    ([5,2,9,7,4,1], 0, 5, '1 е минимален → разменяме с позиция 0'),
    ([1,2,9,7,4,5], 1, 1, '2 е минимален (вече е на място)'),
    ([1,2,9,7,4,5], 2, 4, '4 е минимален → разменяме с позиция 2'),
    ([1,2,4,7,9,5], 3, 5, '5 е минимален → разменяме с позиция 3'),
    ([1,2,4,5,9,7], 4, 5, '7 е минимален → разменяме с позиция 4'),
]

cell_w = 1.15
x_start = 0.4

for row_i, (arr, pos, min_idx, desc) in enumerate(steps_sel):
    y = 2.45 + row_i * 0.88
    for ci, val in enumerate(arr):
        is_pos = ci == pos
        is_min = ci == min_idx
        is_done = ci < pos
        if is_done:
            bg = RGBColor(0xC8, 0xE6, 0xC9)
        elif is_pos:
            bg = C_BLUE
        elif is_min:
            bg = C_ORANGE
        else:
            bg = RGBColor(0xEE, 0xEE, 0xEE)
        fg = C_WHITE if (is_pos or is_min) else (C_GREEN if is_done else C_DARK)
        rect(sl, x_start + ci * cell_w, y, cell_w - 0.05, 0.72, fill=bg)
        txb(sl, str(val), x_start + ci * cell_w + 0.35, y + 0.12,
            0.45, 0.5, size=22, bold=True, color=fg, align=PP_ALIGN.CENTER)
    txb(sl, f'▸  {desc}', x_start + 6.2 * cell_w + 0.1, y + 0.12,
        5.8, 0.6, size=17, color=C_DARK)

# Легенда
rect(sl, 0.4, 7.0, 1.0, 0.36, fill=C_BLUE)
txb(sl, '  текуща позиция', 1.45, 7.0, 3.5, 0.36, size=16, color=C_DARK)
rect(sl, 5.0, 7.0, 1.0, 0.36, fill=C_ORANGE)
txb(sl, '  минимален', 6.05, 7.0, 3.0, 0.36, size=16, color=C_DARK)
rect(sl, 9.2, 7.0, 1.0, 0.36, fill=RGBColor(0xC8,0xE6,0xC9))
txb(sl, '  готов', 10.25, 7.0, 2.5, 0.36, size=16, color=C_DARK)

add_notes(sl, """SELECTION SORT — ИДЕЯ

Разлика от Bubble Sort:
- Bubble Sort: разменя СЪСЕДНИ елементи при всяка стъпка → много размени
- Selection Sort: прави САМО ЕДНА размяна на обход (минималния на правилното място) → по-малко размени

Обяснение стъпка по стъпка:
Обход 1 (i=0): Намираме минималния от цял масив → 1 (на позиция 5)
              Разменяме го с позиция 0: [1, 2, 9, 7, 4, 5]

Обход 2 (i=1): Намираме минималния от позиция 1 нататък → 2 (вече е на позиция 1)
              Не се налага размяна: [1, 2, 9, 7, 4, 5]

И така нататък...

Аналогия: "Представете си картите. Намирате най-малката, слагате я на първо място. После намирате следващата най-малка от останалите и я слагате на второ място. Повтаряте."

Брой обходи: n-1 (точно като Bubble Sort)
Брой размени: максимум n-1 (по-малко от Bubble Sort!)""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 7 — Selection Sort: код
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8,0xF8,0xFF))
header(sl, '{ }  Selection Sort — Код в C#', 'Реализация на метода на пряката селекция')

sel_code = """static void SelectionSort(int[] numbers)
{
    for (int i = 0; i < numbers.Length - 1; i++)
    {
        int min   = numbers[i];   // приемаме текущия за минимален
        int index = i;            // запомняме индекса му

        for (int j = i + 1; j < numbers.Length; j++)
        {
            if (numbers[j] < min)   // намерихме по-малък?
            {
                min   = numbers[j]; // обновяваме минималния
                index = j;          // запомняме новия индекс
            }
        }

        // Разменяме минималния с текущата позиция
        int temp       = numbers[i];
        numbers[i]     = numbers[index];
        numbers[index] = temp;
    }
}"""

code_box(sl, sel_code, 0.35, 1.75, 8.3, 5.5, size=15)

rect(sl, 8.85, 1.75, 4.15, 5.5, fill=C_WHITE, line_color=C_ORANGE)
txb(sl, '🔑  Ключови моменти', 9.05, 1.9, 3.85, 0.5, size=21, bold=True, color=C_ORANGE)
mtxb(sl, [
    ('▸  min и index запомват', False, C_GRAY),
    ('    НАЙ-МАЛКИЯ елемент', False, C_GRAY),
    ('    и неговия ИНДЕКС', False, C_GRAY),
    ('', False, C_GRAY),
    ('▸  Вътрешният цикъл', False, C_GRAY),
    ('    започва от i+1', False, C_GRAY),
    ('    (не проверяваме вече', False, C_GRAY),
    ('    наредените!)', False, C_GRAY),
    ('', False, C_GRAY),
    ('▸  Размяна само веднъж', False, C_GRAY),
    ('    на обход — с temp', False, C_GRAY),
    ('', False, C_GRAY),
    ('▸  За НАМАЛЯВАЩ ред:', False, C_GRAY),
    ('    min → max', False, C_GRAY),
    ('    < → >', False, C_GRAY),
], 9.05, 2.45, 3.85, 4.6, size=18, color=C_GRAY, spacing=5)

add_notes(sl, """SELECTION SORT — КОД

Обяснение на кода:

int min = numbers[i];
int index = i;
→ Приемаме, че ТЕКУЩИЯТ елемент е минималният.
→ Запомняме и ИНДЕКСА му (нужен за размяната!)

Вътрешният цикъл: for (int j = i + 1; ...)
→ Търсим по-малък от i+1 нататък
→ if (numbers[j] < min) → намерихме по-малък → обновяваме min и index

Размяната накрая:
→ Разменяме numbers[i] с numbers[index] (минималния)
→ Само ЕДНА размяна на обход!

Честа грешка: Да разменим с numbers[min] вместо numbers[index]
→ min е СТОЙНОСТТА, index е ПОЗИЦИЯТА!

Сравнение Bubble vs Selection:
Bubble Sort: много размени, лесен за разбиране
Selection Sort: малко размени, малко по-сложен
Двата имат еднаква сложност O(n²) — бавни за голям n""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 8 — Сортиране на низове
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8,0xF8,0xFF))
header(sl, '🔤  Сортиране на масив от низове', 'Не може > или < — използваме CompareTo()')

txb(sl, 'Низовете не могат да се сравняват с > или <  !  Използваме метода CompareTo()',
    0.35, 1.75, 12.6, 0.55, size=22, bold=True, color=C_DARK)

rect(sl, 0.35, 2.4, 12.6, 1.1, fill=RGBColor(0xFF,0xF0,0xE0), line_color=C_ORANGE)
txb(sl, '  низ1.CompareTo(низ2)  връща:', 0.5, 2.48, 5, 0.42, size=20, bold=True, color=C_ORANGE)
mtxb(sl, [
    ('  < 0  →  низ1 е ПРЕДИ низ2 в азбучен ред  (напр. "apple".CompareTo("banana") < 0)',
     False, RGBColor(0x00,0x7A,0x36)),
    ('  = 0  →  низовете са РАВНИ',
     False, C_BLUE),
    ('  > 0  →  низ1 е СЛЕД низ2 в азбучен ред   (напр. "lemon".CompareTo("apple") > 0)',
     False, RGBColor(0xB7,0x1C,0x1C)),
], 0.5, 2.9, 12.1, 0.55, size=18, color=C_GRAY, spacing=4)

str_code = """// Задача: Въвеждаме текст от малки латински букви,
// извеждаме думите и дължините им в азбучен ред.

string[] words = Console.ReadLine().Split();

// Bubble Sort за низове — само if условието се сменя!
for (int i = 0; i < words.Length - 1; i++)
{
    for (int j = 0; j < words.Length - 1 - i; j++)
    {
        if (words[j].CompareTo(words[j + 1]) > 0)  // ← разлика!
        {
            string temp   = words[j];
            words[j]      = words[j + 1];
            words[j + 1]  = temp;
        }
    }
}

foreach (var word in words)
{
    Console.WriteLine($"{word} - {word.Length}");
}"""

code_box(sl, str_code, 0.35, 3.6, 9.0, 3.65, size=14)

rect(sl, 9.55, 3.6, 3.45, 3.65, fill=C_WHITE, line_color=C_BLUE)
txb(sl, '📥  Вход:', 9.75, 3.72, 3.1, 0.4, size=19, bold=True, color=C_BLUE)
txb(sl, 'banana lemon apple', 9.75, 4.15, 3.1, 0.42, size=18, color=C_DARK)
txb(sl, '📤  Изход:', 9.75, 4.7, 3.1, 0.4, size=19, bold=True, color=C_ORANGE)
txb(sl, 'apple - 5\nbanana - 6\nlemon - 5', 9.75, 5.13, 3.1, 1.0, size=18, color=C_DARK)

add_notes(sl, """СОРТИРАНЕ НА НИЗОВЕ

Защо не можем да използваме > и < за низове?
В C# операторите > и < не са дефинирани за тип string.
Ако напишем: if (words[j] > words[j+1]) → ГРЕШКА при компилация!

CompareTo() — как работи?
"apple".CompareTo("banana") → връща отрицателно число (apple е преди banana азбучно)
"banana".CompareTo("apple") → връща положително число
"apple".CompareTo("apple") → връща 0

Важно: CompareTo() е ЧУВСТВИТЕЛЕН към главни/малки букви!
"Apple".CompareTo("apple") → може да даде неочакван резултат
→ За сравнение без значение на главни: използвайте .ToLower() или StringComparison.OrdinalIgnoreCase

Кодът е ИДЕНТИЧЕН с Bubble Sort за числа — само:
- Типът е string вместо int
- if условието е .CompareTo() > 0 вместо >

Примерно обяснение:
"Само едно нещо се сменя в кода — условието if!
Вместо numbers[j] > numbers[j+1] пишем words[j].CompareTo(words[j+1]) > 0""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 9 — Сортиране на обекти
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8,0xF8,0xFF))
header(sl, '🎓  Сортиране на масив от обекти', 'Клас Student — сортиране по успех')

obj_code = """class Student
{
    public string Name     { get; set; }
    public double Absences { get; set; }
    public double Grades   { get; set; }

    public override string ToString()
    {
        return $"{Name} - отс.:{Absences:F1}, успех:{Grades:F2}";
    }
}

// В Main — Bubble Sort по поле Grades (успех):
for (int i = 0; i < n - 1; i++)
{
    for (int j = 0; j < n - 1 - i; j++)
    {
        if (students[j].Grades < students[j + 1].Grades) // ← намаляващ по успех
        {
            Student s       = students[j];
            students[j]     = students[j + 1];
            students[j + 1] = s;
        }
    }
}"""

code_box(sl, obj_code, 0.35, 1.75, 8.5, 5.5, size=14)

rect(sl, 9.1, 1.75, 3.9, 5.5, fill=C_WHITE, line_color=C_ORANGE)
txb(sl, '🔑  Ключова разлика', 9.3, 1.88, 3.6, 0.5, size=20, bold=True, color=C_ORANGE)
mtxb(sl, [
    ('Сравняваме полето', False, C_GRAY),
    ('на обекта:', False, C_GRAY),
    ('', False, C_GRAY),
    ('students[j].Grades', True, C_BLUE),
    ('  вместо', False, C_GRAY),
    ('numbers[j]', True, C_BLUE),
    ('', False, C_GRAY),
    ('Разменяме ЦЕЛИЯ', False, C_GRAY),
    ('обект, не само', False, C_GRAY),
    ('полето!', True, C_ORANGE),
    ('', False, C_GRAY),
    ('Student s = students[j];', False, RGBColor(0x1A,0x5F,0x00)),
    ('students[j] = students[j+1];', False, RGBColor(0x1A,0x5F,0x00)),
    ('students[j+1] = s;', False, RGBColor(0x1A,0x5F,0x00)),
], 9.3, 2.45, 3.65, 4.6, size=16, color=C_GRAY, spacing=5)

add_notes(sl, """СОРТИРАНЕ НА ОБЕКТИ

Ключово разбиране:
При сортиране на обекти:
1. СРАВНЯВАМЕ по конкретно поле (критерий за сортиране)
2. РАЗМЕНЯМЕ ЦЕЛИЯ ОБЕКТ (не само полето!)

Пример за грешка:
Ако разменим само grades:
students[j].Grades = students[j+1].Grades; → ГРЕШКА! Обектите не се разместват, само едно поле се сменя!

Правилно: Разменяме целия Student обект с временна променлива.

Може да сортираме по РАЗЛИЧНИ критерии:
- По успех (Grades) — примерът горе
- По отсъствия (Absences)
- По имена (Name) — с CompareTo()
- По няколко критерия едновременно (по-сложно)

ToString() метода:
public override string ToString() → предефинираме ToString() за красив изход
Console.WriteLine(student) → автоматично вика ToString()
→ Резултат: "Иван Иванов - отс.:12.0, успех:5.20"

Входни данни за задачата:
Брой ученици: 3
Иван Иванов → 12 → 5.2
Петър Петров → 20 → 6
Димитър Димитров → 10 → 4.8""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 10 — Array.Sort() — вградени методи
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0xF8,0xF8,0xFF))
header(sl, '⚡  Array.Sort() — Вградени методи', 'Бързо и лесно сортиране от стандартната библиотека')

txb(sl, 'Не е нужно да пишем сами алгоритъма — C# има вградени методи!',
    0.35, 1.75, 12.6, 0.52, size=23, bold=True, color=C_DARK)

builtin_code = """// 1. Сортиране на числов масив (възходящ ред)
int[] numbers = Console.ReadLine().Split().Select(int.Parse).ToArray();
Array.Sort(numbers);
Console.WriteLine(string.Join(" ", numbers));

// 2. Сортиране в НАМАЛЯВАЩ ред
Array.Sort(numbers);
Array.Reverse(numbers);   // ← обръщаме след сортиране

// 3. Сортиране на масив от низове
string[] words = Console.ReadLine().Split();
Array.Sort(words);        // ← работи директно!

// 4. Сортиране на обекти — трябва IComparable!
// Класът трябва да имплементира CompareTo():
class Student : IComparable<Student>
{
    public int CompareTo(Student other)
    {
        if (Grades < other.Grades) return  1;  // намаляващ
        if (Grades > other.Grades) return -1;
        else                       return  0;
    }
}
// После:
Array.Sort(students);  // ← автоматично използва CompareTo()"""

code_box(sl, builtin_code, 0.35, 2.38, 9.5, 4.9, size=14)

rect(sl, 10.1, 2.38, 2.9, 4.9, fill=C_WHITE, line_color=C_GREEN)
txb(sl, '✅  Предимства', 10.28, 2.52, 2.6, 0.45, size=20, bold=True, color=C_GREEN)
mtxb(sl, [
    ('▸  Бърз (Quick Sort)', False, C_GRAY),
    ('▸  Вграден в C#', False, C_GRAY),
    ('▸  Малко код', False, C_GRAY),
    ('▸  За числа и', False, C_GRAY),
    ('    низове — директно', False, C_GRAY),
    ('', False, C_GRAY),
    ('⚠️  За обекти:', False, C_ORANGE),
    ('Трябва да', False, C_GRAY),
    ('имплементираме', False, C_GRAY),
    ('IComparable!', True, C_ORANGE),
], 10.28, 3.05, 2.65, 4.05, size=17, color=C_GRAY, spacing=6)

add_notes(sl, """ARRAY.SORT() — ВГРАДЕНИ МЕТОДИ

Защо да учим Bubble и Selection Sort ако има Array.Sort()?
"Array.Sort() е като калкулатор — полезен, но трябва да знаете как работи ръчното сортиране, за да разберете какво прави компютърът."
→ На изпит може да се иска ръчно написан алгоритъм!
→ В реалния код → Array.Sort()

Array.Sort() използва вътрешно Quick Sort / Introsort — O(n log n)!

За НАМАЛЯВАЩ ред:
Двете стъпки: Array.Sort() + Array.Reverse()
ВАЖНО: Array.Reverse() не сортира! То просто обръща реда на елементите.
Затова ПЪРВО сортираме, ПОСЛЕ обръщаме.

IComparable:
За обекти Array.Sort() не знае по кое поле да сортира.
Трябва да "кажем" на C# как да сравнява:
→ Класът имплементира IComparable<T>
→ Дефинираме метод CompareTo(), който казва кой е "по-малък"
→ return 1: this е ПО-ГОЛЯМ от other (идва след него)
→ return -1: this е ПО-МАЛЪК от other (идва преди него)

За НАМАЛЯВАЩ по успех: if (Grades < other.Grades) return 1 (по-малкият успех е "по-голям")""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 11 — Сравнение на алгоритмите
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=C_BG)
header(sl, '📊  Сравнение на алгоритмите', 'Кога да ползваме кой?')

headers_t = ['', 'Bubble Sort', 'Selection Sort', 'Array.Sort()']
col_x = [0.35, 2.7, 5.6, 8.55]
col_w = [2.25, 2.8, 2.85, 4.45]
rows_t = [
    ('Сложност', 'O(n²)', 'O(n²)', 'O(n log n)'),
    ('Брой размени', 'Много', 'Малко (n-1)', 'Оптимален'),
    ('Лесен за разбиране', '✅ Да', '✅ Да', '✅ (само извикване)'),
    ('За числа', '✅', '✅', '✅ Array.Sort(a)'),
    ('За низове', '✅ с CompareTo()', '✅ с CompareTo()', '✅ Array.Sort(a)'),
    ('За обекти', '✅ по поле', '✅ по поле', '✅ с IComparable'),
    ('Кога да ползваме', 'Учене, малки масиви', 'Учене, малки масиви', 'В реален код'),
]

y0 = 1.75
rect(sl, 0.35, y0, 12.6, 0.5, fill=C_DARK)
for c, (h, x, w) in enumerate(zip(headers_t, col_x, col_w)):
    txb(sl, h, x+0.05, y0+0.07, w-0.1, 0.36, size=18, bold=True, color=C_WHITE)

for i, (label, b, s, a) in enumerate(rows_t):
    bg = RGBColor(0xEE,0xEE,0xFF) if i % 2 == 0 else C_WHITE
    y = y0 + 0.5 + i * 0.68
    for c, (val, x, w) in enumerate(zip([label,b,s,a], col_x, col_w)):
        rect(sl, x, y, w, 0.68, fill=bg if c > 0 else RGBColor(0xDD,0xDD,0xEE))
        col = C_DARK if c == 0 else (C_GREEN if '✅' in val else C_GRAY)
        txb(sl, val, x+0.08, y+0.1, w-0.16, 0.5, size=16,
            bold=(c==0), color=col)

add_notes(sl, """СРАВНЕНИЕ НА АЛГОРИТМИТЕ

Ключови изводи за учениците:
1. Bubble Sort и Selection Sort са УЧЕБНИ алгоритми — за разбиране на принципа
2. В реален C# код → Array.Sort() (бърз, вграден, проверен)
3. За изпит → може да се иска ръчно написан алгоритъм

Кога Bubble Sort е по-добър от Selection Sort?
Оптимизиран Bubble Sort (с флаг за размяна) може да спре РАНО ако масивът вече е сортиран → O(n) в най-добрия случай!
Selection Sort винаги прави n-1 обхода → O(n²) дори за сортиран масив.

Практически съвет:
"За задачи в клас и изпити → пишете Bubble Sort или Selection Sort.
За проекти и реален код → Array.Sort().
Ако ви питат защо Array.Sort() е по-бързо → защото използва Quick Sort, а не O(n²) алгоритъм."

Въпрос за дискусия: "Ако имате 10 студента → кой алгоритъм? Ако имате 1 000 000 студента → кой?"
→ Малко: няма значение. Много: Array.Sort() задължително!""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 12 — Задача за решаване в клас
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=C_BG)
header(sl, '✏️  Задача за решаване в клас', 'Приложете наученото!')

rect(sl, 0.35, 1.75, 12.6, 2.45, fill=RGBColor(0xFF,0xF0,0xE0), line_color=C_ORANGE)
txb(sl, '📋  Задача:', 0.55, 1.88, 3, 0.5, size=24, bold=True, color=C_ORANGE)
mtxb(sl, [
    ('Напишете програма, която въвежда произволен брой реални числа, разделени с интервал,', False, C_DARK),
    ('и извежда ТРИТЕ НАЙ-ГОЛЕМИ числа в редицата.', True, C_ORANGE),
], 0.55, 2.4, 12.1, 1.6, size=21, color=C_DARK, spacing=10)

rect(sl, 0.35, 4.35, 5.9, 2.95, fill=C_WHITE, line_color=C_BLUE)
txb(sl, '💡  Подход:', 0.55, 4.48, 5.5, 0.45, size=21, bold=True, color=C_BLUE)
mtxb(sl, [
    ('1.  Въведи числата в масив double[]', False, None),
    ('2.  Сортирай масива в НАМАЛЯВАЩ ред', False, None),
    ('3.  Изведи a[0], a[1], a[2]', False, None),
    ('', False, C_GRAY),
    ('Hint: Array.Sort() + Array.Reverse()', True, C_BLUE),
    ('  или Bubble Sort с < вместо >', True, C_BLUE),
], 0.55, 4.98, 5.5, 2.15, size=20, color=C_GRAY, spacing=9)

rect(sl, 6.6, 4.35, 6.1, 2.95, fill=C_WHITE, line_color=C_GREEN)
txb(sl, '✅  Примерен вход/изход:', 6.8, 4.48, 5.7, 0.45, size=21, bold=True, color=C_GREEN)
txb(sl, 'Вход:   3.5  1.2  9.8  4.6  7.1  2.3', 6.8, 4.98, 5.7, 0.5, size=19, color=C_DARK)
txb(sl, 'Изход:  9.8  7.1  4.6', 6.8, 5.55, 5.7, 0.5, size=19, bold=True, color=C_GREEN)
txb(sl, '(трите най-големи, в намаляващ ред)', 6.8, 6.08, 5.7, 0.45, size=17, color=C_GRAY, italic=True)

add_notes(sl, """ЗАДАЧА ЗА КЛАС

Очаквано решение:
double[] numbers = Console.ReadLine()
    .Split()
    .Select(double.Parse)
    .ToArray();

Array.Sort(numbers);
Array.Reverse(numbers);

Console.WriteLine($"{numbers[0]} {numbers[1]} {numbers[2]}");

Или с Bubble Sort в намаляващ ред (< вместо >):
for (int i = 0; i < numbers.Length - 1; i++)
    for (int j = 0; j < numbers.Length - 1 - i; j++)
        if (numbers[j] < numbers[j + 1])  // ← намаляващ
        { ... размяна ... }

Console.WriteLine($"{numbers[0]} {numbers[1]} {numbers[2]}");

Типични грешки:
1. Забравят Array.Reverse() след Array.Sort()
2. Пишат > вместо < за намаляващ Bubble Sort
3. Четат double вместо int (или обратно)

Разширение за бързите:
"Какво ще стане ако масивът има по-малко от 3 елемента?"
→ Трябва проверка: if (numbers.Length >= 3)""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 13 — Обобщение
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=C_DARK)
rect(sl, 0, 0, 0.18, 7.5, fill=C_ORANGE)
rect(sl, 0, 7.35, 13.33, 0.15, fill=C_ORANGE)

txb(sl, '📌  Обобщение — Урок 10', 0.4, 0.15, 12, 0.75,
    size=34, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

points = [
    ('🫧', 'Bubble Sort: сравняваме СЪСЕДНИ, разменяме при нужда — O(n²)'),
    ('🎯', 'Selection Sort: намираме МИНИМАЛНИЯ, разменяме на правилна позиция — O(n²)'),
    ('🔤', 'Низове: не > или <, а  низ1.CompareTo(низ2) > 0'),
    ('🎓', 'Обекти: сравняваме по ПОЛЕ, разменяме ЦЕЛИЯ обект'),
    ('⚡', 'Array.Sort() — вграден, бърз (O(n log n)), за числа и низове директно'),
    ('🎓', 'За обекти с Array.Sort() → имплементираме IComparable<T>'),
]

for i, (icon, text) in enumerate(points):
    y = 1.05 + i * 1.02
    rect(sl, 0.4, y, 12.5, 0.9, fill=RGBColor(0x22,0x22,0x3A))
    txb(sl, f' {icon}   {text}', 0.55, y+0.1, 12.2, 0.72, size=20, color=C_WHITE)

txb(sl, '▶  Следващ урок: Двумерен масив (матрица)',
    0.4, 7.1, 12.5, 0.35, size=15, color=RGBColor(0x80,0x80,0x90),
    italic=True, align=PP_ALIGN.CENTER)

add_notes(sl, """ОБОБЩЕНИЕ — ФИНАЛЕН СЛАЙД

Въпроси за бърза проверка:

1. "Кое число 'изплува' в края след 1 обход на Bubble Sort?"
   → Най-голямото (при сортиране във възходящ ред)

2. "Колко размени прави Selection Sort на обход?"
   → Точно 1

3. "Как сравняваме два низа в C#?"
   → str1.CompareTo(str2)

4. "Кой вграден метод ни сортира масив?"
   → Array.Sort(масив)

5. "Как да сортираме в намаляващ ред с вградени методи?"
   → Array.Sort() + Array.Reverse()

Домашна работа (по желание):
- Задача 2 от учебника: най-голямото число от същите цифри
- Посетете visualgo.net/en/sorting и разгледайте анимациите
- Опитайте да напишете Insertion Sort (подсказка: как слагате карта в ръката си?)""")

# ── Запис ────────────────────────────────────────────────────────────────────
out = Path(r'C:\Users\Neli Nqgolova\Documents\Education\11б клас\Модул_2_Урок_10_Сортиране_на_масив.pptx')
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(out))
print(f'✅ Презентацията е записана:\n   {out}')
print(f'📊 Слайдове: {len(prs.slides)}')
print(f'📝 Бележки за презентатора: на всеки слайд')
