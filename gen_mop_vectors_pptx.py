import sys
sys.stdout.reconfigure(encoding='utf-8')
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree

# ── Цветова схема ──────────────────────────────────────────────────────────
C_DARK   = RGBColor(0x1A, 0x23, 0x7E)
C_ACCENT = RGBColor(0x00, 0xB0, 0xFF)
C_GREEN  = RGBColor(0x00, 0xC8, 0x53)
C_ORANGE = RGBColor(0xFF, 0x6D, 0x00)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xE8, 0xEA, 0xFF)
C_GRAY   = RGBColor(0x33, 0x33, 0x33)
C_BG     = RGBColor(0xF4, 0xF6, 0xFF)
C_CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Помощни функции ─────────────────────────────────────────────────────────
def new_slide():
    return prs.slides.add_slide(BLANK)

def rect(sl, l, t, w, h, fill=None, line_color=None, line_pt=1.5):
    s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color; s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    return s

def txb(sl, text, l, t, w, h, size=22, bold=False, color=C_GRAY,
        align=PP_ALIGN.LEFT, italic=False):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return tb

def mtxb(sl, lines, l, t, w, h, size=20, color=C_GRAY, spacing=6):
    """Multi-line textbox — lines = list of (text, bold, color_override_or_None)"""
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (txt, bld, clr) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(spacing)
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bld
        r.font.color.rgb = clr if clr else color
    return tb

def header(sl, title, sub=''):
    rect(sl, 0, 0, 13.33, 1.55, fill=C_DARK)
    rect(sl, 0, 1.55, 13.33, 0.07, fill=C_ACCENT)
    txb(sl, title, 0.35, 0.08, 12.6, 0.9, size=34, bold=True, color=C_WHITE)
    if sub:
        txb(sl, sub, 0.35, 0.95, 12.6, 0.55, size=19, color=C_ACCENT)

def code_box(sl, code, l, t, w, h, size=16):
    rect(sl, l, t, w, h, fill=C_CODE_BG, line_color=C_ACCENT)
    txb(sl, code, l+0.18, t+0.15, w-0.36, h-0.3, size=size, color=C_GREEN)

def add_notes(sl, text):
    notes_slide = sl.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 1 — Заглавен
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=C_DARK)
rect(sl, 0, 2.5, 13.33, 0.07, fill=C_ACCENT)
rect(sl, 0, 7.43, 13.33, 0.07, fill=C_ACCENT)

txb(sl, 'МОП  ·  Раздел 5  ·  11А клас', 0.6, 0.5, 12, 0.6,
    size=22, color=C_ACCENT)
txb(sl, 'Връзка между вектори\nи масиви в програмирането', 0.6, 1.15, 12, 2.4,
    size=44, bold=True, color=C_WHITE)
txb(sl, 'C# · Unity · Реални приложения', 0.6, 3.75, 10, 0.65,
    size=24, color=RGBColor(0xB3,0xC5,0xFF))
txb(sl, '← Предишен урок: Вектор. Свойства на векторите', 0.6, 4.55, 11, 0.55,
    size=18, color=RGBColor(0x78,0x90,0xAA), italic=True)
txb(sl, 'Нели Няголова  ·  2025/2026', 0.6, 6.7, 10, 0.55,
    size=16, color=RGBColor(0x60,0x70,0x90))

add_notes(sl, """ЗАГЛАВЕН СЛАЙД
Цел на урока: Учениците да разберат как математическият вектор се реализира като масив в C# и защо тази връзка е фундаментална за програмирането.

Въвеждащ въпрос към класа: "Спомняте ли си от миналия час — какво е вектор? Как ще го 'запишем' в компютъра?"
Изчакайте отговори, после преминете към слайд 2.""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 2 — Преговор: Какво е вектор?
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=C_BG)
header(sl, '🔁 Преговор — Какво е вектор?', 'От предишния урок')

rect(sl, 0.35, 1.75, 12.6, 5.5, fill=C_WHITE, line_color=C_DARK)

mtxb(sl, [
    ('📐  Вектор = насочена отсечка с начало и край', True, C_DARK),
    ('', False, C_GRAY),
    ('▸  Означение:  ā  или  AB⃗', False, None),
    ('▸  Характеристики:  посока  +  дължина (модул)', False, None),
    ('▸  В равнината:  ā = (x, y)  — две координати', False, None),
    ('▸  В пространството:  ā = (x, y, z)  — три координати', False, None),
    ('▸  Нулев вектор:  0̄ = (0, 0)  или  (0, 0, 0)', False, None),
    ('▸  Единичен вектор (верзор):  |ā| = 1', False, None),
], 0.65, 2.0, 12.1, 5.0, size=24, color=C_GRAY, spacing=8)

add_notes(sl, """ПРЕГОВОР — Какво е вектор?

Теория за презентатора:
Векторът е математически обект, определен от ПОСОКА и ДЪЛЖИНА (за разлика от точката, която има само позиция, и скалара, който е само число).

Ключови понятия:
- Координати: в 2D — (x, y), в 3D — (x, y, z)
- Модул (дължина): |ā| = √(x² + y²) за 2D
- Нулев вектор: всички координати са 0
- Верзор: единичен вектор, показва само посоката

Пример за клас: "Ако вятърът духа на изток с 5 м/с, векторът на вятъра е (5, 0). Ако духа на североизток — (3.5, 3.5)."

Въпрос: "Как ще съхраним вектора (3, -1, 5) в C#?" → водете към следващия слайд.""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 3 — Преговор: Свойства на векторите
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=C_BG)
header(sl, '🔁 Преговор — Свойства на векторите', 'Правилата, по които работят')

rect(sl, 0.35, 1.75, 6.1, 5.5, fill=C_WHITE, line_color=C_DARK)
txb(sl, '➕  Аритметични свойства', 0.55, 1.9, 5.7, 0.5, size=22, bold=True, color=C_DARK)
mtxb(sl, [
    ('▸  Събиране:  ā + b̄ = (x₁+x₂, y₁+y₂)', False, None),
    ('▸  Комутативност:  ā + b̄ = b̄ + ā', False, None),
    ('▸  Асоциативност:  (ā+b̄)+c̄ = ā+(b̄+c̄)', False, None),
    ('▸  Умножение по скалар:  k·ā = (kx, ky)', False, None),
    ('▸  Нулев елемент:  ā + 0̄ = ā', False, None),
], 0.55, 2.45, 5.7, 4.6, size=22, color=C_GRAY, spacing=10)

rect(sl, 6.9, 1.75, 6.1, 5.5, fill=C_WHITE, line_color=C_ORANGE)
txb(sl, '🔢  Специални операции', 7.1, 1.9, 5.7, 0.5, size=22, bold=True, color=C_ORANGE)
mtxb(sl, [
    ('▸  Модул:  |ā| = √(x²+y²+z²)', False, None),
    ('▸  Скаларно произведение:', False, None),
    ('    ā·b̄ = x₁x₂ + y₁y₂ + z₁z₂', False, RGBColor(0x1A,0x5E,0x00)),
    ('▸  Ако  ā·b̄ = 0  →  ā ⊥ b̄', False, None),
    ('▸  Ъгъл:  cos θ = ā·b̄ / (|ā|·|b̄|)', False, None),
], 7.1, 2.45, 5.7, 4.6, size=22, color=C_GRAY, spacing=10)

add_notes(sl, """СВОЙСТВА НА ВЕКТОРИТЕ

Теория за презентатора:
Тези свойства са МАТЕМАТИЧЕСКАТА ОСНОВА на всичко, което ще правим в C#.

Важни случаи за скаларното произведение:
- ā·b̄ = 0  →  векторите са ПЕРПЕНДИКУЛЯРНИ
- ā·b̄ > 0  →  ъгълът между тях е остър (< 90°)
- ā·b̄ < 0  →  ъгълът между тях е тъп (> 90°)

Пример:
ā = (1, 0),  b̄ = (0, 1)  →  ā·b̄ = 1·0 + 0·1 = 0  →  перпендикулярни (90°) ✔

Модул: ā = (3, 4)  →  |ā| = √(9+16) = √25 = 5

Свържете с C# в следващия слайд: "Всяко от тези действия ще се превърне в for цикъл!".""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 4 — Голямата идея: векторът = масив
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=C_BG)
header(sl, '💡 Голямата идея', 'Векторът е масив — масивът е вектор!')

rect(sl, 0.35, 1.75, 12.6, 1.0, fill=C_DARK)
txb(sl, '  ā = (3,  -1,  5)   ←→   double[] a = { 3, -1, 5 };',
    0.5, 1.88, 12.3, 0.72, size=28, bold=True, color=C_WHITE)

rect(sl, 0.35, 2.95, 5.9, 4.2, fill=C_WHITE, line_color=C_DARK)
txb(sl, '📐  Математика', 0.55, 3.1, 5.5, 0.5, size=22, bold=True, color=C_DARK)
mtxb(sl, [
    ('Вектор:  ā = (x₁, x₂, ..., xₙ)', False, None),
    ('Размерност:  n', False, None),
    ('Координата:  xᵢ  (i = 1..n)', False, None),
    ('Нулев вектор:  (0, 0, ..., 0)', False, None),
], 0.55, 3.65, 5.5, 3.3, size=22, color=C_GRAY, spacing=12)

txb(sl, '⟺', 6.4, 4.5, 0.9, 0.7, size=38, bold=True,
    color=C_ORANGE, align=PP_ALIGN.CENTER)

rect(sl, 7.45, 2.95, 5.55, 4.2, fill=C_WHITE, line_color=C_ACCENT)
txb(sl, '{ }  C# масив', 7.65, 3.1, 5.15, 0.5, size=22, bold=True, color=C_ACCENT)
mtxb(sl, [
    ('double[] a = new double[n];', False, C_GREEN),
    ('Дължина:  a.Length', False, C_GRAY),
    ('Елемент:  a[i]  (i = 0..n-1)', False, C_GRAY),
    ('Нули:  new double[n]  по подразбиране', False, C_GRAY),
], 7.65, 3.65, 5.15, 3.3, size=22, color=C_GRAY, spacing=12)

add_notes(sl, """ГОЛЯМАТА ИДЕЯ

Теория за презентатора:
Това е ЦЕНТРАЛНАТА ИДЕЯ на урока — векторът от математиката и масивът в програмирането са едно и също нещо, просто записани по различен начин.

Математиката използва индексиране от 1: x₁, x₂, x₃
C# използва индексиране от 0: a[0], a[1], a[2]
→ Това е важна разлика! Напомнете го.

Пример за обяснение:
"Векторът ā = (3, -1, 5) е тридimensионен — има 3 координати.
В C# го пишем: double[] a = { 3, -1, 5 };
a[0] = 3   (това е x-координатата)
a[1] = -1  (това е y-координатата)
a[2] = 5   (това е z-координатата)"

Въпрос към класа: "Ако имаме 10-мерен вектор, как ще го декларираме в C#?"
Отговор: double[] a = new double[10];""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 5 — Таблица с паралели
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=C_BG)
header(sl, '📊 Паралел: математика ↔ C#', 'Всяко понятие има еквивалент в кода')

rows = [
    ('Координата  xᵢ',             'a[i]',                          'Индексиране от 0 в C#!'),
    ('Размерност  n',               'a.Length',                      'Брой елементи'),
    ('Нулев вектор  (0,0,...)',     'new double[n]',                 'Нули по подразбиране'),
    ('Сумиране  ā + b̄',             'c[i] = a[i] + b[i]  (for)',     'Елемент по елемент'),
    ('Умножение  k·ā',              'b[i] = k * a[i]  (for)',        'Скаларно умножение'),
    ('Скаларно произведение',       'sum += a[i]*b[i]  (for)',       'Резултатът е число'),
    ('Модул  |ā|',                  'Math.Sqrt(Σ a[i]*a[i])',        'Питагорова теорема'),
]

y0 = 1.72
rect(sl, 0.3, y0, 4.5, 0.52, fill=C_DARK)
rect(sl, 4.85, y0, 4.3, 0.52, fill=C_DARK)
rect(sl, 9.2, y0, 3.85, 0.52, fill=C_DARK)
txb(sl, '  Математика', 0.3, y0+0.07, 4.4, 0.38, size=19, bold=True, color=C_WHITE)
txb(sl, '  C# код', 4.85, y0+0.07, 4.2, 0.38, size=19, bold=True, color=C_WHITE)
txb(sl, '  Бележка', 9.2, y0+0.07, 3.75, 0.38, size=19, bold=True, color=C_WHITE)

for i, (m, c, n) in enumerate(rows):
    bg = RGBColor(0xEE,0xF0,0xFF) if i % 2 == 0 else C_WHITE
    rh = 0.62
    y = y0 + 0.52 + i * rh
    rect(sl, 0.3, y, 4.5, rh, fill=bg)
    rect(sl, 4.85, y, 4.3, rh, fill=bg)
    rect(sl, 9.2, y, 3.85, rh, fill=bg)
    txb(sl, '  ' + m, 0.3, y+0.08, 4.4, rh-0.1, size=18, color=C_DARK)
    txb(sl, '  ' + c, 4.85, y+0.08, 4.2, rh-0.1, size=18, color=RGBColor(0x00,0x6B,0x00), bold=True)
    txb(sl, '  ' + n, 9.2, y+0.08, 3.75, rh-0.1, size=16, color=C_GRAY, italic=True)

add_notes(sl, """ТАБЛИЦА С ПАРАЛЕЛИ

Използвайте тази таблица като справочник — можете да се върнете към нея по всяко време на урока.

Важни акценти:
1. Индексиране от 0: В математиката x₁ е първата координата, в C# a[0] е първият елемент.
2. for цикълът е универсалният инструмент за операции с вектори.
3. Скаларното произведение дава ЧИСЛО (не вектор!) — това е честа грешка при учениците.

Пример за скаларно произведение:
ā = (1, 2, 3),  b̄ = (4, 5, 6)
ā·b̄ = 1·4 + 2·5 + 3·6 = 4 + 10 + 18 = 32

Пример за модул:
ā = (3, 4)  →  |ā| = √(9+16) = 5""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 6 — C# код: Събиране и умножение
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0x10,0x10,0x1A))
header(sl, '{ }  C# — Събиране и умножение по скалар', 'Операциите от математиката като код')

code_add = """// Събиране на два вектора
static double[] Add(double[] a, double[] b)
{
    double[] c = new double[a.Length];
    for (int i = 0; i < a.Length; i++)
        c[i] = a[i] + b[i];
    return c;
}"""

code_scale = """// Умножение на вектор по скалар
static double[] Scale(double[] a, double k)
{
    double[] result = new double[a.Length];
    for (int i = 0; i < a.Length; i++)
        result[i] = a[i] * k;
    return result;
}"""

code_box(sl, code_add,   0.35, 1.75, 6.1, 4.0, size=17)
code_box(sl, code_scale, 6.6,  1.75, 6.4, 4.0, size=17)

rect(sl, 0.35, 5.9, 12.65, 1.35, fill=RGBColor(0x1A,0x23,0x7E))
txb(sl, '▸  double[] a = {1,2,3};   double[] b = {4,5,6};',
    0.55, 5.97, 12.2, 0.5, size=18, color=C_WHITE)
txb(sl, '▸  Add(a,b)   →  {5, 7, 9}       Scale(a, 2)   →  {2, 4, 6}',
    0.55, 6.52, 12.2, 0.5, size=18, color=C_GREEN, bold=True)

add_notes(sl, """C# КОД: СЪБИРАНЕ И УМНОЖЕНИЕ

Теория за презентатора:
И двете операции следват един и същ МОДЕЛ:
1. Създаваме нов масив с размера на входните
2. Обхождаме с for цикъл всеки индекс
3. Прилагаме операцията елемент по елемент
4. Връщаме резултата

ВАЖНО — типична грешка на учениците:
double[] c = a;  // ГРЕШНО! c и a сочат към един и същ масив!
double[] c = new double[a.Length];  // ПРАВИЛНО!

Демонстрация на дъската:
ā = (1, 2, 3)   b̄ = (4, 5, 6)
i=0: c[0] = a[0]+b[0] = 1+4 = 5
i=1: c[1] = a[1]+b[1] = 2+5 = 7
i=2: c[2] = a[2]+b[2] = 3+6 = 9
Резултат: c = (5, 7, 9) ✔

Scale(a, 2): умножаваме всеки елемент по 2 → (2, 4, 6)""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 7 — C# код: Модул и скаларно произведение
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0x10,0x10,0x1A))
header(sl, '{ }  C# — Модул и скаларно произведение', 'Операции, чийто резултат е ЧИСЛО')

code_mag = """// Модул (дължина) на вектор
static double Magnitude(double[] a)
{
    double sum = 0;
    for (int i = 0; i < a.Length; i++)
        sum += a[i] * a[i];
    return Math.Sqrt(sum);
}"""

code_dot = """// Скаларно произведение
static double DotProduct(double[] a, double[] b)
{
    double sum = 0;
    for (int i = 0; i < a.Length; i++)
        sum += a[i] * b[i];
    return sum;
}"""

code_box(sl, code_mag, 0.35, 1.75, 6.1, 3.8, size=17)
code_box(sl, code_dot, 6.6,  1.75, 6.4, 3.8, size=17)

rect(sl, 0.35, 5.7, 12.65, 1.55, fill=RGBColor(0x1A,0x23,0x7E))
txb(sl, '▸  Magnitude({3, 4})   →   5.0   (защото √(9+16) = √25 = 5)',
    0.55, 5.78, 12.2, 0.55, size=18, color=C_WHITE)
txb(sl, '▸  DotProduct({1,0}, {0,1})   →   0.0   →   ПЕРПЕНДИКУЛЯРНИ! ⊥',
    0.55, 6.35, 12.2, 0.55, size=18, color=C_GREEN, bold=True)

add_notes(sl, """МОДУЛ И СКАЛАРНО ПРОИЗВЕДЕНИЕ

Теория за презентатора:
Модулът (Magnitude) дава ДЪЛЖИНАТА на вектора:
|ā| = √(x₁² + x₂² + ... + xₙ²)
Това е просто Питагоровата теорема в n-мерно пространство!

Скаларното произведение (Dot Product):
ā·b̄ = x₁y₁ + x₂y₂ + ... + xₙyₙ

Критично важно приложение:
Ако DotProduct(a, b) == 0, то a ⊥ b (перпендикулярни)
→ Това се използва НАВСЯКЪДЕ в графиката и физиката на игрите!

Пример:
ā = (3, 4, 0):  Magnitude = √(9+16+0) = √25 = 5
ā = (1, 0, 0),  b̄ = (0, 1, 0):  DotProduct = 0  →  оси x и y са ⊥

Въпрос: "Как да проверим дали два вектора са перпендикулярни?"
Отговор: DotProduct(a, b) == 0 (или Math.Abs(...) < 0.0001 за реални числа)""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 8 — Нормализиране: единичен вектор
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0x10,0x10,0x1A))
header(sl, '{ }  C# — Нормализиране (единичен вектор)', 'Запазваме посоката, правим дължина = 1')

code_norm = """// Нормализиране на вектор (верзор)
// Резултатът има дължина = 1, но същата посока
static double[] Normalize(double[] a)
{
    double mag = Magnitude(a);

    if (mag == 0)
        throw new Exception("Не може да се нормализира нулев вектор!");

    double[] result = new double[a.Length];
    for (int i = 0; i < a.Length; i++)
        result[i] = a[i] / mag;

    return result;
}"""

code_box(sl, code_norm, 0.35, 1.75, 7.5, 5.45, size=17)

rect(sl, 8.1, 1.75, 4.9, 2.6, fill=RGBColor(0x1A,0x23,0x7E))
txb(sl, '🔑  Защо е нужно?', 8.3, 1.88, 4.55, 0.45, size=21, bold=True, color=C_ACCENT)
mtxb(sl, [
    ('▸  Искаме само ПОСОКАТА,', False, C_WHITE),
    ('    не дължината', False, C_WHITE),
    ('▸  Използва се в:', False, C_WHITE),
    ('    игри, физика, ML', False, C_ACCENT),
], 8.3, 2.4, 4.55, 2.7, size=19, color=C_WHITE, spacing=8)

rect(sl, 8.1, 4.5, 4.9, 2.7, fill=C_CODE_BG, line_color=C_GREEN)
txb(sl, '// Пример:\ndouble[] v = {3, 4, 0};\n// mag = 5\n// result = {0.6, 0.8, 0}\n// Magnitude(result) = 1 ✔',
    8.28, 4.65, 4.55, 2.4, size=16, color=C_GREEN)

add_notes(sl, """НОРМАЛИЗИРАНЕ

Теория за презентатора:
Нормализирането (normalization) превръща вектора в единичен вектор (верзор):
â = ā / |ā|

Резултатът:
- Има СЪЩАТА ПОСОКА като оригиналния вектор
- Има ДЪЛЖИНА = 1

Пример:
ā = (3, 4, 0)
|ā| = 5
â = (3/5, 4/5, 0) = (0.6, 0.8, 0)
|â| = √(0.36 + 0.64 + 0) = √1 = 1 ✔

Реално приложение (Unity):
В игрите, когато движим персонаж:
Vector3 dir = new Vector3(dx, 0, dz).normalized;
transform.position += dir * speed * Time.deltaTime;
Без нормализиране, движението по диагонал е по-бързо (√2 пъти)!

Грешка: Никога не нормализирайте нулевия вектор (деление на нула)!""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 9 — Защо ги учат? Реални приложения
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=C_BG)
header(sl, '🎯 Защо учим вектори в програмирането?', 'Вектори са навсякъде около нас!')

apps = [
    ('🎮', 'Игри (Unity / C#)',
     '• Позиция, скорост, посока на персонажа\n• Vector3 pos = new Vector3(x, y, z)\n• Всяко движение = вектор!',
     C_DARK),
    ('🤖', 'Изкуствен интелект',
     '• Думите в ChatGPT са вектори от 1536 числа!\n• "цар" - "мъж" + "жена" ≈ "царица"\n• ML.NET, TensorFlow работят с масиви',
     RGBColor(0x5A,0x00,0x8A)),
    ('🗺️', 'GPS и навигация',
     '• Посоката на движение = вектор\n• Разстояние = модул на вектора\n• Google Maps изчислява вектори',
     RGBColor(0x00,0x6B,0x3F)),
    ('📊', 'Анализ на данни',
     '• Всеки запис в база данни = вектор\n• Сходство = скаларно произведение\n• Spotify, Netflix препоръчват с вектори',
     RGBColor(0xB7,0x1C,0x1C)),
]

for i, (icon, title, desc, col) in enumerate(apps):
    x = 0.35 + (i % 2) * 6.5
    y = 1.75 + (i // 2) * 2.8
    rect(sl, x, y, 6.15, 2.6, fill=C_WHITE, line_color=col)
    txb(sl, f'{icon}  {title}', x+0.2, y+0.12, 5.75, 0.5, size=22, bold=True, color=col)
    txb(sl, desc, x+0.2, y+0.68, 5.75, 1.8, size=18, color=C_GRAY)

add_notes(sl, """ЗАЩО УЧИМ ВЕКТОРИ?

Цел на слайда: Мотивация! Учениците трябва да разберат, че това НЕ е абстрактна математика.

Разкажете историята с ChatGPT (много впечатлява!):
"ChatGPT представя всяка дума като вектор от 1536 числа. Когато GPT 'разбира' текст, той всъщност извършва операции с вектори — ТОЧНО КАТО ТЕ ПРАВЯТ ДНЕС!
Сходството между думи се изчислява с... скаларно произведение!"

За Unity:
"Ако искате да правите игри с Unity — 90% от кода работи с Vector2 и Vector3. Всичко, което учите днес, е директно приложимо."

За Netflix/Spotify:
"Когато Spotify ви препоръчва песен, той изчислява скаларното произведение между вектора на вашите предпочитания и векторите на песните."

Въпрос: "Кой от тези примери ви е най-интересен?"
→ Нека учениците споделят.""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 10 — Unity: Vector3
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0x10,0x10,0x1A))
header(sl, '🎮  Unity + C# — Vector3 в игрите', 'Математическият вектор оживява!')

unity_code = """public class Player : MonoBehaviour
{
    public float speed = 5f;

    void Update()
    {
        // Вектор на движение от клавиатурата
        Vector3 dir = new Vector3(
            Input.GetAxis("Horizontal"),  // x
            0f,                           // y (без скачане)
            Input.GetAxis("Vertical")     // z
        );

        // Нормализираме — еднаква скорост
        if (dir.magnitude > 0)
            dir = dir.normalized;

        // Движим се: позиция += посока * скорост * време
        transform.position += dir * speed * Time.deltaTime;
    }
}"""

code_box(sl, unity_code, 0.35, 1.75, 8.0, 5.5, size=15)

rect(sl, 8.55, 1.75, 4.45, 5.5, fill=RGBColor(0x1A,0x23,0x7E))
txb(sl, '🔑  Vector3 методи', 8.75, 1.88, 4.1, 0.5, size=21, bold=True, color=C_ACCENT)
mtxb(sl, [
    ('.magnitude  →  модул', False, C_WHITE),
    ('.normalized  →  верзор', False, C_WHITE),
    ('Vector3.Dot(a,b)  →  ā·b̄', False, C_WHITE),
    ('Vector3.Cross(a,b)  →  a×b', False, C_WHITE),
    ('Vector3.Distance(a,b)', False, C_WHITE),
    ('Vector3.Lerp(a,b,t)', False, C_WHITE),
    ('Vector3.zero  =  (0,0,0)', False, C_WHITE),
    ('Vector3.up  =  (0,1,0)', False, C_WHITE),
], 8.75, 2.45, 4.1, 4.55, size=18, color=C_WHITE, spacing=10)

add_notes(sl, """UNITY + C# — Vector3

Теория за презентатора:
Unity е безплатен game engine, използван от МНОГО компании. C# е основният му език. Vector3 е вградена структура, която е точно 3D вектор.

Обяснение на кода:
1. Input.GetAxis("Horizontal") → стойност -1 до 1 (клавиши A/D или стрелки)
2. dir = нов вектор (вляво/вдясно, нагоре/надолу, напред/назад)
3. dir.normalized → единичен вектор (без него движението по диагонал е по-бързо!)
4. transform.position += dir * speed * Time.deltaTime
   → позицията се мести в посоката на вектора
   → Time.deltaTime гарантира еднаква скорост на всички компютри

Връзка с урока:
transform.position е Vector3 = (x, y, z) = double[3]
dir.magnitude е точно Magnitude() функцията от предния слайд!
dir.normalized е точно Normalize() функцията от предния слайд!

Кажете: "Кодът, който написахте в C# Console App, е ИДЕНТИЧЕН с кода в Unity!".""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 11 — Дигитални игри и ресурси
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=C_BG)
header(sl, '🎮  Интерактивни ресурси за учене', 'Опитайте ги у дома или в клас!')

resources = [
    ('🟦', 'GeoGebra', 'geogebra.org',
     'Визуализирайте вектори в 2D и 3D.\nСъбирайте, въртете, измервайте — интерактивно!',
     C_DARK),
    ('🟡', 'Kahoot', 'kahoot.com',
     'Търсете "vectors" или "вектори" —\nнамерете готови викторини за класа!',
     C_ORANGE),
    ('🟢', 'Unity Learn', 'learn.unity.com',
     '"Create with Code" — безплатен курс.\nПравите 3D игра с Vector3 в C#!',
     RGBColor(0x00,0x7A,0x36)),
    ('🔴', 'Khan Academy', 'khanacademy.org',
     'Видео уроци по вектори с автопроверка.\nИма и на български!',
     RGBColor(0xB7,0x1C,0x1C)),
    ('🟣', 'Brilliant.org', 'brilliant.org',
     'Интерактивен курс по линейна алгебра.\nСтъпка по стъпка с визуализации.',
     RGBColor(0x6A,0x1B,0x9A)),
    ('🩵', 'Scratch', 'scratch.mit.edu',
     'Анимация с движещи се обекти —\nвекторите оживяват визуално!',
     RGBColor(0x00,0x70,0xA0)),
]

for i, (dot, name, url, desc, col) in enumerate(resources):
    x = 0.35 + (i % 3) * 4.32
    y = 1.75 + (i // 3) * 2.8
    rect(sl, x, y, 4.0, 2.6, fill=C_WHITE, line_color=col)
    txb(sl, f'{dot}  {name}', x+0.18, y+0.12, 3.65, 0.5, size=22, bold=True, color=col)
    txb(sl, url, x+0.18, y+0.62, 3.65, 0.3, size=14, color=C_ACCENT, italic=True)
    txb(sl, desc, x+0.18, y+0.98, 3.65, 1.5, size=17, color=C_GRAY)

add_notes(sl, """ДИГИТАЛНИ ИГРИ И РЕСУРСИ

Препоръки за ползване в клас:

GeoGebra (СИЛНО препоръчително!):
- Отидете на geogebra.org/classic
- Нарисувайте два вектора, покажете тяхната сума графично
- Пресметнете скаларното произведение и покажете геометричния смисъл
- Безплатно, работи в браузъра, без регистрация

Kahoot в клас:
- Създайте 5-6 въпроса за следващия час (10 минути)
- Или намерете готова викторина по "vectors"
- Учениците играят с телефоните си

Unity Learn:
- Препоръчайте за домашна работа
- "Create with Code" е безплатен и с награди (сертификат!)
- Учениците, интересуващи се от игри, ще го харесат

Khan Academy:
- Има СТРАХОТНИ видеа по линейна алгебра
- Препоръчайте за допълнително изучаване""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 12 — Упражнение: Описание
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=C_BG)
header(sl, '✏️  Упражнение — Задача за клас', 'Напишете методите сами!')

rect(sl, 0.35, 1.75, 12.6, 3.7, fill=RGBColor(0xFF,0xF3,0xE0), line_color=C_ORANGE)
txb(sl, '📋  Задача:', 0.55, 1.88, 3, 0.5, size=24, bold=True, color=C_ORANGE)
mtxb(sl, [
    ('Напишете C# метод  VectorInfo(double[] a, double[] b)  който:', True, C_DARK),
    ('', False, C_GRAY),
    ('  1.  Принтира сумата на двата вектора', False, None),
    ('  2.  Принтира скаларното им произведение', False, None),
    ('  3.  Принтира дължината (модула) на първия вектор', False, None),
    ('  4.  Казва дали векторите са перпендикулярни (dot product = 0)', False, None),
], 0.55, 2.4, 12.1, 2.8, size=21, color=C_GRAY, spacing=8)

rect(sl, 0.35, 5.6, 12.6, 1.65, fill=RGBColor(0xE8,0xF5,0xE9), line_color=RGBColor(0x2E,0x7D,0x32))
txb(sl, '✅  Тест с тези данни:', 0.55, 5.72, 6, 0.45, size=21, bold=True, color=RGBColor(0x1B,0x5E,0x20))
txb(sl, 'a = { 3, 4, 0 }    b = { -4, 3, 0 }',
    0.55, 6.18, 6, 0.45, size=20, color=C_GRAY)
txb(sl, 'Сума: {-1, 7, 0}   |  Dot: 0   |  |a|: 5   |  Перпендикулярни: ДА ✔',
    6.8, 6.18, 6.0, 0.45, size=19, color=RGBColor(0x1B,0x5E,0x20), bold=True)

add_notes(sl, """УПРАЖНЕНИЕ — ОПИСАНИЕ

Организация:
- Индивидуална работа: 10-12 минути
- После: преглед на решенията заедно

Насоки ако учениците заседнат:
"Имате нужда от 4 метода — Add(), DotProduct(), Magnitude() — вече ги написахте заедно! Просто ги извикайте от VectorInfo."

Проверка на перпендикулярност:
bool isPerp = Math.Abs(DotProduct(a, b)) < 0.0001;
(Не == 0, защото реалните числа имат грешка при закръгление!)

Разширена задача (за бързите):
Добавете метод Angle(a, b), който изчислява ъгъла между двата вектора:
double angle = Math.Acos(DotProduct(a,b) / (Magnitude(a) * Magnitude(b)));
double degrees = angle * 180 / Math.PI;

Очакван изход за тест данните:
Сума: [-1, 7, 0]
Скаларно произведение: 0
Модул на a: 5
Перпендикулярни: ДА""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 13 — Упражнение: Шаблон
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0x10,0x10,0x1A))
header(sl, '{ }  Шаблон за упражнението', 'Копирайте и попълнете!')

template = """using System;

class VectorOps
{
    static double[] Add(double[] a, double[] b) { /* ... */ }
    static double DotProduct(double[] a, double[] b) { /* ... */ }
    static double Magnitude(double[] a) { /* ... */ }

    static void VectorInfo(double[] a, double[] b)
    {
        double[] sum    = Add(a, b);
        double   dot    = DotProduct(a, b);
        double   mag    = Magnitude(a);
        bool     isPerp = Math.Abs(dot) < 0.0001;

        Console.Write("Сума: [");
        for (int i = 0; i < sum.Length; i++)
            Console.Write((i > 0 ? ", " : "") + sum[i]);
        Console.WriteLine("]");

        Console.WriteLine($"Скаларно произведение: {dot}");
        Console.WriteLine($"Модул на a: {mag:F4}");
        Console.WriteLine($"Перпендикулярни: {(isPerp ? "ДА" : "НЕ")}");
    }

    static void Main()
    {
        double[] a = { 3, 4, 0 };
        double[] b = { -4, 3, 0 };
        VectorInfo(a, b);
    }
}"""

code_box(sl, template, 0.35, 1.75, 12.6, 5.5, size=14)

add_notes(sl, """ШАБЛОН ЗА УПРАЖНЕНИЕТО

Показвайте този слайд СЛЕД като учениците са опитали сами.

Обяснение на шаблона:
1. Трите метода (Add, DotProduct, Magnitude) са НЕПОПЪЛНЕНИ — учениците ги пишат.
2. VectorInfo() е вече написан — показва как да се използват методите.
3. Math.Abs(dot) < 0.0001 вместо dot == 0 — за реални числа!

Console.Write vs Console.WriteLine:
- Write: принтира без нов ред
- WriteLine: принтира с нов ред
→ Използваме Write в for цикъла, за да принтираме масива на един ред.

Форматиране:
$"Модул: {mag:F4}" → принтира с 4 десетични знака
Пример: 5.0000

Разговор след упражнението:
"Забелязахте ли, че Add(), DotProduct() и Magnitude() следват ЕДИН И СЪЩ МОДЕЛ?
for цикъл + операция + return резултат"
→ Това е СИЛАТА на масивите!""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 14 — Разширение: LINQ и foreach
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=RGBColor(0x10,0x10,0x1A))
header(sl, '{ }  Бонус — По-елегантен C# код', 'LINQ и foreach за по-кратък запис')

code_linq = """// С LINQ (Language Integrated Query)
using System.Linq;

// Модул
double mag = Math.Sqrt(a.Select(x => x*x).Sum());

// Скаларно произведение
double dot = a.Zip(b, (x, y) => x * y).Sum();

// Сумиране на вектори
double[] sum = a.Zip(b, (x, y) => x + y).ToArray();

// Проверка за перпендикулярност
bool isPerp = Math.Abs(dot) < 0.0001;"""

code_foreach = """// С foreach (само за четене)
double sumOfSquares = 0;
foreach (double x in a)
    sumOfSquares += x * x;

double magnitude = Math.Sqrt(sumOfSquares);

// Принтиране на вектор
Console.Write("[");
Console.Write(string.Join(", ", a));
Console.WriteLine("]");
// Резултат: [3, -1, 5]"""

code_box(sl, code_linq,    0.35, 1.75, 6.3, 4.8, size=16)
code_box(sl, code_foreach, 6.75, 1.75, 6.25, 4.8, size=16)

rect(sl, 0.35, 6.7, 12.6, 0.6, fill=RGBColor(0x1A,0x23,0x7E))
txb(sl, '💡  И двата начина дават СЪЩИЯ резултат — изберете по-четимия за вас!',
    0.55, 6.78, 12.1, 0.42, size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

add_notes(sl, """LINQ И FOREACH — БОНУС СЛАЙД

Покажете само ако времето позволява или за бързите ученици.

Обяснение на LINQ:
LINQ = Language Integrated Query — позволява операции с колекции като едно изречение.
a.Select(x => x*x) → трансформира всеки елемент (lambda функция)
.Sum() → сумира всички елементи
a.Zip(b, (x,y) => x*y) → комбинира два масива елемент по елемент

Кога да използваме for vs LINQ vs foreach:
- for: когато ни трябва индексът i
- foreach: когато само четем елементите
- LINQ: когато искаме кратък, функционален запис

string.Join(", ", a):
Превръща масива в стринг: [3, -1, 5] → "3, -1, 5"
Много полезно за принтиране!

Това е НАПРЕДНАЛ материал — не е задължителен за упражнението.""")

# ════════════════════════════════════════════════════════════════════════════
# СЛАЙД 15 — Обобщение
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
rect(sl, 0, 0, 13.33, 7.5, fill=C_DARK)
rect(sl, 0, 0, 13.33, 0.07, fill=C_ACCENT)
rect(sl, 0, 7.43, 13.33, 0.07, fill=C_ACCENT)

txb(sl, '📌  Обобщение на урока', 0.5, 0.2, 12, 0.75,
    size=34, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

points = [
    ('🔢', 'Векторът  (x₁, x₂, ..., xₙ)  =  масив  double[]  в C#'),
    ('🔁', 'Всяка векторна операция = for цикъл елемент по елемент'),
    ('📏', 'Модул = Math.Sqrt(Σ a[i]²)     |     Dot product = Σ a[i]·b[i]'),
    ('⊥',  'Перпендикулярни вектори: DotProduct(a, b) == 0'),
    ('🎮', 'Unity Vector3 е директна реализация на 3D вектор в C#'),
    ('🌍', 'Приложения: игри, AI, GPS, анализ на данни — навсякъде!'),
]

for i, (icon, text) in enumerate(points):
    y = 1.05 + i * 1.02
    rect(sl, 0.5, y, 12.3, 0.9, fill=RGBColor(0x22,0x2D,0x8E))
    txb(sl, f' {icon}   {text}', 0.65, y+0.1, 12.0, 0.7, size=21, color=C_WHITE)

txb(sl, '▶  Следващ урок: Операции с вектори — задачи и приложения',
    0.5, 7.1, 12.3, 0.35, size=15, color=RGBColor(0x80,0x90,0xB0),
    italic=True, align=PP_ALIGN.CENTER)

add_notes(sl, """ОБОБЩЕНИЕ — ФИНАЛЕН СЛАЙД

Затваряне на урока (5 минути):

Въпроси за проверка на разбирането:
1. "Как декларираме 4D вектор ā = (1, 2, 3, 4) в C#?"
   → double[] a = { 1, 2, 3, 4 };

2. "Кой индекс е третата координата на вектора?"
   → a[2]  (индексиране от 0!)

3. "Два вектора са перпендикулярни. Какво е скаларното им произведение?"
   → 0

4. "Как изчисляваме модула?"
   → Math.Sqrt(a[0]*a[0] + a[1]*a[1] + ...) или с Magnitude() метода

Домашна работа (по желание):
- Напишете метод Angle(a, b) — ъгъл между два вектора в градуси
- Отидете на geogebra.org и нарисувайте сумата на два вектора
- Разгледайте Unity Learn: "Create with Code"

Следващ урок: Задачи с вектори — прилагаме всичко наученото.""")

# ── Запис ────────────────────────────────────────────────────────────────────
out = Path(r'C:\Users\Neli Nqgolova\Documents\Education\11а клас\МОП\Раздел 5\Вектори_и_масиви_C#.pptx')
prs.save(str(out))
print(f'✅ Презентацията е записана:\n   {out}')
print(f'📊 Слайдове: {len(prs.slides)}')
print(f'📝 Всеки слайд съдържа бележки за презентатора')
